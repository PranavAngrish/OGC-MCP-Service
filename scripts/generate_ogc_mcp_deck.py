from pathlib import Path
import sys

sys.path.insert(0, "/private/tmp/ogc-pptx-deps")

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parents[1] / "OGC_to_MCP_Bridging_Pranav_Angrish.pptx"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(5, 20, 18)
PANEL = RGBColor(12, 39, 35)
PANEL_2 = RGBColor(16, 52, 46)
GREEN = RGBColor(143, 215, 93)
MINT = RGBColor(100, 226, 194)
BLUE = RGBColor(101, 174, 241)
AMBER = RGBColor(248, 192, 86)
WHITE = RGBColor(239, 247, 241)
MUTED = RGBColor(163, 190, 179)
RED = RGBColor(245, 123, 119)


def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height)
    shape.fill.solid(); shape.fill.fore_color.rgb = GREEN; shape.line.fill.background()
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.12), 0, prs.slide_width, Inches(0.05))
    shape.fill.solid(); shape.fill.fore_color.rgb = PANEL_2; shape.line.fill.background()


def text(slide, value, x, y, w, h, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = value
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return box


def title(slide, heading, kicker="OGC API MCP BRIDGE"):
    text(slide, kicker, 0.55, 0.30, 5.2, 0.28, 10, GREEN, True)
    text(slide, heading, 0.55, 0.65, 12.1, 0.58, 29, WHITE, True)


def footer(slide, number):
    text(slide, f"{number:02d}", 12.42, 7.02, 0.45, 0.22, 10, MUTED, True, PP_ALIGN.RIGHT)
    text(slide, "Pranav Angrish  ·  GSoC 2026", 0.55, 7.02, 3.0, 0.22, 9, MUTED)


def card(slide, x, y, w, h, heading, body="", accent=GREEN, heading_size=15):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = PANEL; s.line.color.rgb = PANEL_2
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    strip.fill.solid(); strip.fill.fore_color.rgb = accent; strip.line.fill.background()
    text(slide, heading, x + 0.23, y + 0.18, w - 0.42, 0.33, heading_size, WHITE, True)
    if body:
        text(slide, body, x + 0.23, y + 0.58, w - 0.42, h - 0.72, 11.5, MUTED)
    return s


def compact_card(slide, x, y, w, h, heading, detail="", accent=GREEN, heading_size=11, detail_size=8.5):
    """A dense card for architecture diagrams with a short label and detail."""
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = PANEL; s.line.color.rgb = PANEL_2
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    strip.fill.solid(); strip.fill.fore_color.rgb = accent; strip.line.fill.background()
    text(slide, heading, x + 0.15, y + 0.09, w - 0.27, 0.22, heading_size, WHITE, True)
    if detail:
        text(slide, detail, x + 0.15, y + 0.34, w - 0.27, h - 0.39, detail_size, MUTED)
    return s


def architecture_band(slide, y, h, label, accent):
    """Create a readable horizontal layer for the architecture slide."""
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.83), Inches(y), Inches(10.98), Inches(h))
    band.fill.solid(); band.fill.fore_color.rgb = PANEL; band.line.color.rgb = PANEL_2
    marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.52), Inches(y), Inches(1.05), Inches(h))
    marker.fill.solid(); marker.fill.fore_color.rgb = accent; marker.line.color.rgb = accent
    text(slide, label, 0.58, y + h * 0.31, 0.93, h * 0.30, 8.2, BG, True, PP_ALIGN.CENTER)
    return band


def arrow(slide, x1, y1, x2, y2, color=MINT, width=2.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color; line.line.width = Pt(width); line.line.end_arrowhead = True
    return line


def circle(slide, x, y, d, label, fill, label_size=13):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.color.rgb = fill
    text(slide, label, x, y + d * .27, d, d * .35, label_size, BG, True, PP_ALIGN.CENTER)
    return s


# 1 — title
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
text(slide, "GSoC 2026  ·  52°North", 0.72, 0.72, 4.0, 0.3, 12, GREEN, True)
text(slide, "OGC to MCP\nBridging", 0.72, 1.25, 7.1, 1.55, 42, WHITE, True)
text(slide, "Making geospatial APIs usable through safe,\nauditable natural-language workflows.", 0.76, 3.15, 5.8, 0.72, 18, MUTED)
text(slide, "by Pranav Angrish", 0.76, 5.86, 3.5, 0.32, 16, MINT, True)
circle(slide, 8.4, 1.45, 1.5, "OGC\nAPIs", GREEN)
circle(slide, 10.55, 1.45, 1.5, "MCP", BLUE)
circle(slide, 9.48, 4.05, 1.5, "AI +\nHuman", AMBER)
arrow(slide, 9.9, 2.2, 10.52, 2.2); arrow(slide, 11.0, 2.92, 10.45, 4.0); arrow(slide, 9.45, 4.0, 8.98, 2.92)
text(slide, "trusted bridge", 9.36, 3.12, 1.7, 0.26, 10, MUTED, False, PP_ALIGN.CENTER)
footer(slide, 1)

# 2 — problem
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "The problem: one question hides a chain of expert-only decisions")
text(slide, "A normal user sees one task. An OGC API workflow exposes many coupled technical choices — and one wrong choice can invalidate the analysis.", 0.56, 1.32, 12.0, 0.36, 14, MUTED)
card(slide, 0.62, 2.05, 2.3, 2.75, "What a user asks", "“Find suitable point data, run a spatial analysis, and show me the result on a map.”", GREEN, 16)
text(slide, "one natural-language request", 0.72, 4.47, 2.1, 0.24, 10, GREEN, True, PP_ALIGN.CENTER)
arrow(slide, 3.02, 3.42, 3.55, 3.42, MINT, 2.5)
text(slide, "becomes", 3.08, 3.02, 0.42, 0.22, 10, MUTED, True, PP_ALIGN.CENTER)
steps = [
    ("Discover", "server + conformance"),
    ("Understand", "collections + schema"),
    ("Query", "filters + CRS + paging"),
    ("Execute", "process inputs + jobs"),
    ("Interpret", "formats + output validity"),
]
for i, (head, sub) in enumerate(steps):
    y = 1.95 + i * 0.76
    c = [GREEN, BLUE, AMBER, RED, MINT][i]
    circle(slide, 3.72, y, 0.42, str(i + 1), c, 10)
    card(slide, 4.28, y - 0.06, 3.0, 0.54, head, sub, c, 12)
arrow(slide, 7.47, 3.42, 8.0, 3.42, RED, 2.5)
text(slide, "or risk", 7.47, 3.02, 0.7, 0.22, 10, RED, True, PP_ALIGN.CENTER)
card(slide, 8.18, 1.88, 3.85, 1.0, "Wrong / incomplete evidence", "The answer may be based on the wrong collection, filter scope, CRS, or truncated result.", RED, 14)
card(slide, 8.18, 3.17, 3.85, 1.0, "Unsafe process execution", "An incorrect input, remote reference, or unreviewed request can trigger an unintended operation.", RED, 14)
card(slide, 8.18, 4.46, 3.85, 1.0, "Output that cannot be trusted", "Raw GeoJSON/GML, asynchronous jobs, and unknown formats do not automatically become a usable map.", RED, 14)
text(slide, "The gap is not just convenience — it is correctness, safety, and accessibility.", 1.35, 6.08, 10.5, 0.35, 17, MINT, True, PP_ALIGN.CENTER)
footer(slide, 2)

# 3 — solution
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "Proposed solution: turn intent into a verified OGC workflow")
text(slide, "The OGC-to-MCP bridge adds the missing decision, safety, and presentation layers between a user and fragmented geospatial APIs.", 0.56, 1.32, 12.0, 0.34, 14, MUTED)
card(slide, 0.56, 1.92, 2.22, 3.62, "User intent", "“Find suitable data, run this analysis, and show me the result.”", GREEN, 16)
text(slide, "Natural language", 0.77, 5.10, 1.75, 0.22, 10, GREEN, True, PP_ALIGN.CENTER)
arrow(slide, 2.84, 3.66, 3.20, 3.66, MINT, 2.5)

solution_steps = [
    ("1", "Discover safely", "Uses only registered servers; reads conformance and capabilities.", GREEN),
    ("2", "Build the right query", "Inspects collections and schemas; validates filters, scope, CRS, and paging.", BLUE),
    ("3", "Create a reviewable plan", "Constructs one exact process request; validates input schema and sources.", AMBER),
    ("4", "Require human approval", "Shows the exact request and fingerprint before any state-changing execution.", RED),
    ("5", "Verify and present", "Monitors jobs; normalizes outputs into trustworthy maps, tables, and downloads.", MINT),
]
solution_positions = [(3.28, 1.86), (6.33, 1.86), (9.38, 1.86), (9.38, 3.94), (6.33, 3.94)]
for (number, heading, body, color), (x, y) in zip(solution_steps, solution_positions):
    card(slide, x, y, 2.78, 1.58, heading, body, color, 13)
    circle(slide, x + 2.37, y - 0.17, 0.45, number, color, 10)
arrow(slide, 6.10, 2.65, 6.30, 2.65, MINT, 1.7)
arrow(slide, 9.15, 2.65, 9.35, 2.65, MINT, 1.7)
arrow(slide, 10.77, 3.47, 10.77, 3.90, MINT, 1.7)
arrow(slide, 9.35, 4.73, 9.15, 4.73, MINT, 1.7)

card(slide, 3.28, 6.05, 8.87, 0.55, "What the user receives", "An auditable answer: validated evidence + exact approved execution + verified output artifact.", MINT, 13)
footer(slide, 3)

# 4 — architecture
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "How one request becomes a verified geospatial result")
text(slide, "Three layers have distinct responsibilities: understand the user, control the workflow, and safely work with OGC data.", 0.56, 1.32, 11.9, 0.30, 13.5, MUTED)

# First layer: conversation
architecture_band(slide, 1.84, 1.00, "1  CONVERSATION", GREEN)
compact_card(slide, 2.03, 2.04, 2.55, 0.60, "User + React workspace", "Ask a question · review an exact plan · receive maps and outputs", GREEN, 11.5, 8.3)
compact_card(slide, 5.02, 2.04, 2.55, 0.60, "Gemini LLM", "Turns intent into a small set of proposed MCP tool calls", BLUE, 11.5, 8.3)
compact_card(slide, 9.00, 2.04, 3.00, 0.60, "What comes back", "Answer · activity timeline · approval request · verified visualization", MINT, 11.5, 8.3)
arrow(slide, 4.63, 2.34, 4.97, 2.34, GREEN, 2.0)
arrow(slide, 7.63, 2.34, 8.92, 2.34, MINT, 1.6)
text(slide, "intent", 4.63, 2.58, 0.34, 0.16, 8, MUTED, False, PP_ALIGN.CENTER)

# Second layer: control
architecture_band(slide, 3.16, 1.26, "2  CONTROL", BLUE)
compact_card(slide, 2.03, 3.36, 3.00, 0.86, "Node gateway", "Executes allowed tools · streams SSE updates\nMonitors asynchronous jobs · scopes session artifacts", BLUE, 11.5, 8.4)
compact_card(slide, 5.42, 3.36, 3.00, 0.86, "Human-in-the-loop gate", "The LLM cannot confirm execution.\nA user reviews a fingerprint-bound exact request.", RED, 11.5, 8.4)
compact_card(slide, 8.81, 3.36, 3.20, 0.86, "MCP protocol boundary", "Structured ogc_* tool calls over stdio\nwith stable contracts and result envelopes", AMBER, 11.5, 8.4)
arrow(slide, 6.28, 2.84, 6.28, 3.32, BLUE, 2.0)
text(slide, "only permitted tool requests", 6.48, 2.96, 1.5, 0.16, 8, MUTED)

# Third layer: data proxy
architecture_band(slide, 4.78, 1.22, "3  DATA + GOVERNANCE", AMBER)
compact_card(slide, 2.03, 4.97, 2.48, 0.82, "OGC MCP proxy", "Knows the OGC standards and exposes one predictable tool surface", AMBER, 11.3, 8.3)
compact_card(slide, 4.82, 4.97, 2.48, 0.82, "Validation + security", "Registered servers · schema checks · auth injection · bounded HTTP", RED, 11.3, 8.3)
compact_card(slide, 7.61, 4.97, 2.48, 0.82, "Large-data boundary", "Raw GeoJSON/GML stays in proxy memory behind opaque handles", BLUE, 11.3, 8.3)
compact_card(slide, 10.40, 4.97, 1.61, 0.82, "Artifacts", "Verify\nmap/table\ndownload", MINT, 10.5, 8.0)
arrow(slide, 10.41, 4.19, 10.41, 4.93, AMBER, 2.0)
text(slide, "MCP", 10.55, 4.48, 0.35, 0.16, 8, MUTED)

# External services + foundations
architecture_band(slide, 6.18, 0.53, "OGC API LAYER", MINT)
text(slide, "Operator-approved OGC servers:   Common   ·   Features   ·   Records   ·   Processes   ·   Jobs", 2.10, 6.32, 10.1, 0.18, 10.2, WHITE, True, PP_ALIGN.CENTER)

# High-value boundaries that explain the design at a glance
card(slide, 0.78, 6.76, 1.56, 0.20, "", "", RED, 10)
text(slide, "Human approves execution", 0.83, 6.79, 1.46, 0.10, 6.4, WHITE, True, PP_ALIGN.CENTER)
card(slide, 9.82, 6.76, 2.60, 0.20, "", "", BLUE, 10)
text(slide, "Raw geometry never enters the LLM", 9.87, 6.79, 2.50, 0.10, 6.4, WHITE, True, PP_ALIGN.CENTER)
footer(slide, 4)

# 5 — OGC surface
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "One MCP tool surface over multiple OGC API modules")
text(slide, "Users do not need to memorize endpoints; the bridge selects the correct standard operation.", 0.56, 1.32, 11.5, 0.34, 14, MUTED)
modules = [("Common", "landing · conformance\ncollections", GREEN), ("Features", "collections · items\nqueryables · filters", BLUE), ("Records", "catalogue search\nmetadata records", AMBER), ("Processes", "process schemas\nexecution plans", RED), ("Jobs", "status · results\nasync monitoring", MINT)]
for i, (h, b, c) in enumerate(modules):
    x = 0.58 + i * 2.53
    card(slide, x, 2.1, 2.18, 2.4, h, b, c)
    text(slide, f"ogc_{h.lower()}_*", x + 0.18, 4.12, 1.8, 0.25, 10, c, True, PP_ALIGN.CENTER)
text(slide, "Stable MCP tools", 0.72, 5.25, 2.05, 0.3, 13, GREEN, True, PP_ALIGN.CENTER)
arrow(slide, 2.83, 5.38, 10.45, 5.38)
text(slide, "Standards-aware routing, validation, and fallback", 3.75, 5.68, 5.9, 0.3, 14, WHITE, True, PP_ALIGN.CENTER)
footer(slide, 5)

# 6 — deterministic
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "Making agentic geospatial workflows deterministic")
text(slide, "The model can suggest actions; deterministic services decide what is allowed and what is true.", 0.56, 1.32, 11.4, 0.34, 14, MUTED)
circle(slide, 0.8, 2.55, 1.35, "LLM\nintent", AMBER)
arrow(slide, 2.2, 3.23, 3.02, 3.23)
card(slide, 3.05, 2.05, 2.3, 2.35, "Policy gate", "Only registered servers\nOnly valid tool calls\nNo direct execution", RED)
arrow(slide, 5.42, 3.23, 6.24, 3.23)
card(slide, 6.27, 2.05, 2.3, 2.35, "Validation", "Schema checks\nCapability checks\nBounded requests", BLUE)
arrow(slide, 8.64, 3.23, 9.46, 3.23)
card(slide, 9.49, 2.05, 2.65, 2.35, "Human checkpoint", "Exact request\nFingerprint\nExplicit approval", GREEN)
text(slide, "Deterministic lifecycle", 0.8, 5.2, 2.1, 0.3, 14, MINT, True)
arrow(slide, 2.85, 5.35, 11.35, 5.35, MINT, 2.5)
for x, label in [(3.1, "discover"), (5.15, "plan"), (7.1, "confirm"), (9.25, "execute"), (11.1, "verify")]:
    circle(slide, x, 5.02, 0.58, "", MINT, 1)
    text(slide, label, x - 0.45, 5.75, 1.5, 0.23, 10, MUTED, False, PP_ALIGN.CENTER)
footer(slide, 6)

# 7 — data boundary
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide); title(slide, "Keeping enormous geospatial data out of LLM context")
text(slide, "Raw coordinates remain in the trusted data plane; the model receives only bounded, sanitized facts.", 0.56, 1.32, 11.8, 0.34, 14, MUTED)
card(slide, 0.65, 2.05, 2.4, 3.05, "OGC response", "Large GeoJSON / GML\nMany features\nCoordinate arrays\nRemote references", RED)
arrow(slide, 3.12, 3.55, 3.92, 3.55)
card(slide, 3.95, 2.05, 2.55, 3.05, "Proxy memory", "Full response retained\nOpaque mem_* handle\nTTL-bound storage\nSafe retrieval", BLUE)
arrow(slide, 6.57, 3.55, 7.37, 3.55)
card(slide, 7.4, 2.05, 2.25, 3.05, "LLM context", "Compact summary\nIDs & properties\nNo raw geometry\nNo instructions", AMBER)
arrow(slide, 6.0, 5.55, 6.0, 6.1, MINT)
card(slide, 4.0, 6.1, 4.0, 0.55, "Trusted UI only", "Hydrates verified outputs → map, table, download", GREEN, 13)
text(slide, "Model-safe", 7.58, 5.38, 1.8, 0.25, 11, AMBER, True, PP_ALIGN.CENTER)
text(slide, "Large data never enters the conversational prompt.", 0.72, 5.72, 3.0, 0.32, 14, MINT, True)
footer(slide, 7)

# 8 — thank you
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide)
text(slide, "Thank you", 0.72, 1.30, 6.5, 0.62, 38, WHITE, True)
text(slide, "OGC to MCP Bridging", 0.76, 2.12, 5.3, 0.36, 20, GREEN, True)
text(slide, "Safe. Auditable. Conversational. Geospatial.", 0.76, 2.72, 6.5, 0.34, 16, MUTED)
circle(slide, 8.42, 2.05, 1.55, "OGC", GREEN, 19)
circle(slide, 10.58, 2.05, 1.55, "MCP", BLUE, 19)
arrow(slide, 9.98, 2.83, 10.55, 2.83)
text(slide, "Pranav Angrish", 0.76, 5.73, 3.5, 0.3, 16, MINT, True)
text(slide, "GSoC 2026  ·  52°North", 0.76, 6.16, 3.5, 0.25, 11, MUTED)
footer(slide, 8)

prs.save(OUT)
print(OUT)
